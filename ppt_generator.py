# ppt_generator.py
from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor

def _format_title(shp, text, size=44):
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    try:
        p.font.color.rgb = RGBColor(0, 51, 102)
    except Exception:
        pass

def _add_bullets(slide, bullets, size=20):
    try:
        tf = slide.placeholders[1].text_frame
    except Exception:
        tx = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(3))
        tf = tx.text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(size)
        p.level = 0

def _add_footer(slide, footer_text):
    try:
        left = Cm(1)
        top = Cm(18)
        width = Cm(24)
        height = Cm(0.8)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(10)
        p.font.italic = True
    except Exception:
        pass

def build_ppt_from_slides(slides, presentation_title, output_path="result.pptx", subtitle=None, sources_meta=None):
    """
    slides: ordered list for the six topics (title, bullets, notes)
    """
    prs = Presentation()
    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    _format_title(s0.shapes.title, presentation_title, size=44)
    if subtitle:
        try:
            s0.placeholders[1].text = subtitle
        except Exception:
            pass

    # Content slides (one per slide object)
    for sdata in slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = sdata.get("title", "")
        _add_bullets(s, sdata.get("bullets", []), size=18)
        try:
            s.notes_slide.notes_text_frame.text = sdata.get("notes", "")
        except Exception:
            pass
        if sources_meta:
            _add_footer(s, "Sources included — see Sources slide")

    # Sources slide
    if sources_meta:
        try:
            s_sources = prs.slides.add_slide(prs.slide_layouts[1])
            s_sources.shapes.title.text = "Sources"
            bullets = [f"{m.get('id')}: {m.get('title') or m.get('url')}\n{m.get('url')}" for m in sources_meta]
            _add_bullets(s_sources, bullets, size=12)
        except Exception:
            pass

    # Q&A
    qa = prs.slides.add_slide(prs.slide_layouts[1])
    qa.shapes.title.text = "Q & A"
    _add_bullets(qa, ["Any questions?", "Contact: you@example.com"], size=20)

    prs.save(output_path)
    return output_path